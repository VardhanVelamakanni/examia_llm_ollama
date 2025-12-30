from pptx import Presentation
from PIL import Image
import io
import os

def extract_images_from_ppt(file_path, output_dir="extracted_images"):
    prs = Presentation(file_path)

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    image_count = 0

    for slide_index, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image = Image.open(io.BytesIO(shape.image.blob))
                image_count += 1

                image_path = f"{output_dir}/slide_{slide_index+1}_img_{image_count}.png"
                image.save(image_path)

    return image_count


# TEST
if __name__ == "__main__":
    count = extract_images_from_ppt("DACS-2025_Presentation.pptx")
    print(f"Extracted {count} images")
