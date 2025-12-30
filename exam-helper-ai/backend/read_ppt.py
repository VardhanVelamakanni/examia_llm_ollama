from pptx import Presentation

def read_ppt(file_path):
    prs = Presentation(file_path)
    text = ""

    for slide_number, slide in enumerate(prs.slides, start=1):
        text += f"\n--- Slide {slide_number} ---\n"

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                shape_text = shape.text.strip()
                if shape_text:
                    text += shape_text + "\n"

    return text


# TEST
if __name__ == "__main__":
    content = read_ppt("DACS-2025_Presentation.pptx")
    print(content)
