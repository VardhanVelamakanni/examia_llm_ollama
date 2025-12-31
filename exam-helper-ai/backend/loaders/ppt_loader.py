from pptx import Presentation

def load_ppt(path: str) -> str:
    prs = Presentation(path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text
